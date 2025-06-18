import matplotlib.pyplot as plt
import matplotlib as mpl
from matplotlib import cm
import os
import numpy as np
from qutip import *
from qutip.piqs import *
from qutip import piqs
import PhaseSpaceRepresentations_upgr as ps
from tqdm import tqdm
from pptx import Presentation
from pptx.util import Inches

#import matplotlib.animation as animation
#from IPython.display import HTML
#from IPython.core.display import Image, display


def rotation_y(rho,j_y,N,phi):
    system = piqs.Dicke(N = N)
    system.hamiltonian = 0.5*phi*j_y
    D_tls = system.liouvillian()
    t = np.linspace(0, 1, 100)
    
    result = mesolve(D_tls, rho, t, [])
    return result.states[-1]

def rotation(rho, axis_vector, N, phi):
    """Rotation about arbitrary axis using proper Qobj operators"""
    # Get spin operators
    [jx, jy, jz] = piqs.jspin(N)
    
    # Construct axis operator from vector components
    j_axis = (axis_vector[0]*jx + 
              axis_vector[1]*jy + 
              axis_vector[2]*jz)
    
    # Create Dicke system with rotation Hamiltonian
    system = piqs.Dicke(N=N)
    system.hamiltonian = 0.5 * phi * j_axis  # Now using proper Qobj
    
    # Solve dynamics
    D_tls = system.liouvillian()
    t = np.linspace(0, 1, 100)
    result = mesolve(D_tls, rho, t, [])
    
    return result.states[-1]



def OAT(rho,j_z,N,Q_per_nsc,F_per_nsc,n_sc,scattering=True):
    
    system = piqs.Dicke(N = N)
    t = np.linspace(0, 1, 100) #tau=1
    chi=Q_per_nsc*n_sc/N
    # system.hamiltonian = chi*Qobj(j_z.full() @ j_z.full())
    system.hamiltonian = chi*j_z**2

    if scattering:
        system.dephasing =2*n_sc/N
    system.collective_dephasing=F_per_nsc*n_sc/N
    
    D_tls = system.liouvillian()
    
    result = mesolve(D_tls, rho, t, [])
    return result.states[-1]

def density_op_plot_grid(ax,N):

    lines=[N]
    if N%2==0:
        for S in np.flip(0.5*np.arange(0,N,2)):
            ax.axhline(lines[-1],c='black')
            ax.axvline(lines[-1],c='black')
            lines.append(lines[-1]+2*S+1)    
    else:
        for S in np.flip(0.5*np.arange(1,N,2)):
            ax.axhline(lines[-1],c='black')
            ax.axvline(lines[-1],c='black')
            lines.append(lines[-1]+2*S+1)


def isolate_density_op_block(rho, N, S):
    """
    Extract the block corresponding to total spin S from the density matrix.
    
    Parameters:
    rho: density matrix
    N: number of spins (N=20 in your case)
    S: desired total spin quantum number
    
    Returns:
    block: (2S+1)×(2S+1) density matrix block for spin S
    """
    # Check trace of full density matrix using QuTiP's trace
    # print(f"Shape of rho: {rho.shape}")
    # print(f"Trace of rho: {(rho.tr()).real}")
    # Start from maximum S (N/2) and sum dimensions of all blocks until we reach our desired S
    S_max = N/2
    index = 0
    
    # Sum dimensions of all blocks with higher S values
    for S_block in np.arange(S_max, S, -1):
        index += int(2*S_block + 1)
    
    block_size = int(2*S + 1)
    block = rho[index:index+block_size, index:index+block_size]
    
    # Add verification
    # print(f"For S={S}:")
    # print(f"  Block starts at index {index}")
    # print(f"  Block size: {block_size}×{block_size}")
    # print(f"  Block trace: {np.real(np.trace(block)):.6f}")
    # print(f"  Trace Block^2: {np.real(np.trace(block@block))}")
    
    return block

def calculate_qfi(rho, j_axis):
    return 4 * (expect(j_axis**2, rho) - (expect(j_axis, rho))**2)

def calculate_axis_qfi(rho, axis_vector, jx, jy, jz):
    """Calculate QFI for arbitrary rotation axis"""
    assert np.isclose(np.linalg.norm(axis_vector), 1), "Must be unit vector"
    j_axis = axis_vector[0]*jx + axis_vector[1]*jy + axis_vector[2]*jz
    return 4 * (expect(j_axis**2, rho) - (expect(j_axis, rho))**2)


    # return 4 * (expect(Qobj(j_axis.full() @ j_axis.full()), rho) - (expect(j_axis, rho))**2)
    return 4 * (expect(j_axis**2, rho) - (expect(j_axis, rho))**2)


def random_axis_operator():
    """Generate random unit vector in 3D space"""
    vec = np.random.normal(size=3)
    return vec / np.linalg.norm(vec)



def calculate_qfi_mesh(rho, jx, jy, jz, n_phi=50, n_theta=100, n_samples=10):
    # Create spherical coordinate grid
    phi_vals = np.linspace(0, np.pi, n_phi)
    theta_vals = np.linspace(0, 2*np.pi, n_theta)
    phi_grid, theta_grid = np.meshgrid(phi_vals, theta_vals, indexing='ij')

    # Convert to Cartesian unit vectors
    x = np.sin(phi_grid) * np.cos(theta_grid)
    y = np.sin(phi_grid) * np.sin(theta_grid)
    z = np.cos(phi_grid)
    
    # Initialize storage
    qfi_avg = np.zeros_like(phi_grid)
    qfi_var = np.zeros_like(phi_grid)

    # Calculate QFI for each mesh point
    for i in tqdm(range(n_phi)):
        for j in tqdm(range(n_theta)):
            axis_vec = np.array([x[i,j], y[i,j], z[i,j]])
            j_axis = axis_vec[0]*jx + axis_vec[1]*jy + axis_vec[2]*jz
            qfi = 4 * (expect(j_axis**2, rho) - (expect(j_axis, rho))**2)
            # qfi = 4 * (expect(Qobj(j_axis.full() @ j_axis.full()), rho) - (expect(j_axis, rho))**2)
            qfi_avg[i,j] = qfi
            qfi_var[i,j] = 0  # Variance not needed for single calculation
    
    return phi_grid, theta_grid, qfi_avg, qfi_var

def plot_qfi_mesh(phi, theta, qfi_avg, qfi_var, N, Q_per_nsc, i, outdir='qfi_mesh_results'):
    os.makedirs(outdir, exist_ok=True)
    # Create figure with two subplots
    fig, ax = plt.subplots(1, 2, figsize=(20, 8), 
                         gridspec_kw={'width_ratios': [1.2, 1]})
    
    # Print statistics
    print(f"Maximum QFI: {np.max(qfi_avg):.2f}")
    print(f"Minimum QFI: {np.min(qfi_avg):.2f}")
    
    # Left plot: Polar projection of average QFI
    ax[0] = plt.subplot(121, projection='polar')
    c0 = ax[0].contourf(theta, phi, qfi_avg, 50, cmap='viridis')
    ax[0].set_title('Angular QFI Distribution', pad=20)
    fig.colorbar(c0, ax=ax[0], label='QFI Value')
    ax[0].set_theta_zero_location('N')
    ax[0].set_theta_direction(-1)
    ax[0].set_rlabel_position(90)

    # Right plot: Histogram with statistics
    all_qfi = qfi_avg.flatten()
    avg_qfi = np.mean(all_qfi)
    std_qfi = np.std(all_qfi)
    
    n, bins, patches = ax[1].hist(all_qfi, bins=50, alpha=0.7, 
                                label='QFI Distribution')
    ax[1].axvline(avg_qfi, color='r', linestyle='--', 
                label=f'Average: {avg_qfi:.2f}')
    ax[1].fill_betweenx([0, np.max(n)], 
                      avg_qfi - std_qfi,
                      avg_qfi + std_qfi,
                      color='gray', alpha=0.2, 
                      label=f'±1σ ({std_qfi:.2f})')
    
    ax[1].set_xlabel('Quantum Fisher Information')
    ax[1].set_ylabel('Frequency')
    ax[1].set_title(f'QFI Distribution (N={N}, η={Q_per_nsc})')
    ax[1].legend()
    ax[1].grid(True)

    plt.tight_layout()
    fname = f'iter{i+1}_QFI.png'
    # plt.savefig(fname, bbox_inches='tight')
    plt.show()

    # Print variance
    print(f"Global Variance: {np.var(all_qfi):.2f}")
    
    # filename = os.path.join(outdir, f'qfi_mesh_nsc_{n_sc:.1f}.png')
    # plt.savefig(filename, bbox_inches='tight', dpi=150)
    # plt.close()
    
    # return filename

def create_presentation(filenames):
    """Create PowerPoint with QFI mesh plots"""
    prs = Presentation()
    
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Kicked Top QFI Analysis"
    slide.placeholders[1].text = "Quantum Fisher Information Mesh Plots\nn_sc values: 0.0, 0.1, 0.7"

    # Content slides
    for fname in filenames:
        # Create new slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = os.path.basename(fname).replace('_', ' ').replace('.png', '')
        
        # Add image
        left = Inches(0.5)
        top = Inches(1.5)
        height = Inches(6)
        slide.shapes.add_picture(fname, left, top, height=height)

    prs.save('QFI_Mesh_Analysis.pptx')
    print("PowerPoint presentation created successfully!")


np.random.seed(42)  # Set the seed for reproducibility

N = 28
jp=piqs.jspin(N,"+")
jm=jp.dag()
iters = 100

Q_per_nsc=30#60#0.23 #for N=200
F_per_nsc=0.#9#0.06 #for N=200
n_sc_values=[1]#.7#0.7#0.7#0.1

#From the paper converting k, p to Q and theta, 
# Good region: p = pi/2, k = 6
#k = 2*Q_per_nsc*n_sc/N, 2*6*1
#p = 0.5*angle

n_samples = 1  # Number of axis samples
qfi_values = []
saved_filenames = []

for n_sc in n_sc_values:
    print(f"\n=== Processing n_sc = {n_sc} ===")
    # Generate random axis
    system = piqs.Dicke(N = N)
    jx, jy, jz = piqs.jspin(N)
    axis = random_axis_operator()
    
    # Initialize state
    rho = piqs.dicke(N, N/2, N/2)
    
    # Bring to equator
    # rho = rotation_y(rho,jy,N,np.pi)
    rho = rotation(rho, [0,1,0], N, np.pi)
    
    
    # OAT dynamics
    for i in tqdm(range(iters)):
        rho = OAT(rho, jz, N, Q_per_nsc, F_per_nsc, n_sc=n_sc, scattering=False)
        # axis = random_axis_operator()
        # print(axis)
        angle = np.pi
        rho = rotation(rho, [0,1,0], N, angle)
        # rho = rotation(rho, axis, N, np.pi)

    ### Plotting Wigner function ###
    # ISOLATE SYMMETRIC SUBSPACE (S=14 for N=28)
    S_max = N//2
    rho_block = isolate_density_op_block(rho, N, S_max)
    rho_np = rho_block  # Convert to NumPy array

    # Calculate Wigner function
    Kcoeffs = ps.precalculatedKcoeffs(int(2*S_max + 1), 'Calculated/Kernels/')
    result = ps.PSrepresentationFromFourier(rho_np, Kcoeffs, 384)
    
    # Plot
    plt.figure(figsize=(8,8))
    ps.PSrepPlot_plane(result, f'Wigner Function after {iters} OAT Cycles')
    fname = f'iter{0}_Wigner.png'
    # plt.savefig(fname, bbox_inches='tight')
    plt.show()

    # Simulation parameters
    phi_grid, theta_grid, qfi_avg, qfi_var = calculate_qfi_mesh(rho, jx, jy, jz)
    fname = plot_qfi_mesh(phi_grid, theta_grid, qfi_avg, qfi_var, N, Q_per_nsc, 0)
    saved_filenames.append(fname)

# create_presentation(saved_filenames)
   